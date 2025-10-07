from fastapi import FastAPI, Depends, HTTPException, UploadFile, File
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from pydantic import BaseModel
from typing import Optional
from datetime import datetime, timedelta
from jose import JWTError, jwt
import logging
import fitz  # PDF
import docx  # DOCX
from pptx import Presentation  # PPTX

from services import summarize_text, summarize_with_t5, generate_quiz, extract_keywords, generate_flashcards

# ------------------ Config ------------------
SECRET_KEY = "secretfortest"   # ⚠️ Change in production
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 60
fake_users_db = {}

# ------------------ Logging ------------------
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

# ------------------ FastAPI App ------------------
app = FastAPI(title="Adaptive Quiz & Summarizer Agent (Gemini + T5)")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

# ------------------ Models ------------------
class User(BaseModel):
    username: str
    full_name: Optional[str] = None
    password: str

class Token(BaseModel):
    access_token: str
    token_type: str

class SummaryRequest(BaseModel):
    text: str
    mode: Optional[str] = "short"  # short | detailed

class TextRequest(BaseModel):
    text: str

class QuizRequest(BaseModel):
    text: str
    difficulty: Optional[str] = "medium"  # easy | medium | hard

# ------------------ Utils ------------------
def fake_hash_password(password: str):
    return "hashed-" + password

def verify_password(plain_password, hashed_password):
    return hashed_password == fake_hash_password(plain_password)

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    to_encode = data.copy()
    expire = datetime.utcnow() + (expires_delta or timedelta(minutes=15))
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

async def get_current_user(token: str = Depends(oauth2_scheme)):
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None or username not in fake_users_db:
            raise HTTPException(status_code=401, detail="Invalid credentials")
        return fake_users_db[username]
    except JWTError:
        raise HTTPException(status_code=401, detail="Invalid token")

# ------------------ Auth Routes ------------------
@app.post("/register")
def register(user: User):
    if user.username in fake_users_db:
        raise HTTPException(status_code=400, detail="Username already exists")
    fake_users_db[user.username] = {
        "username": user.username,
        "full_name": user.full_name,
        "hashed_password": fake_hash_password(user.password),
    }
    return {"message": "User registered successfully"}

@app.post("/token", response_model=Token)
async def login(form_data: OAuth2PasswordRequestForm = Depends()):
    user = fake_users_db.get(form_data.username)
    if not user or not verify_password(form_data.password, user["hashed_password"]):
        raise HTTPException(status_code=400, detail="Invalid username or password")
    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    token = create_access_token(data={"sub": user["username"]}, expires_delta=access_token_expires)
    return {"access_token": token, "token_type": "bearer"}

# ------------------ Core Routes ------------------
@app.get("/")
def home():
    return {"message": "Server is running 🚀"}

@app.post("/summarize")
async def summarize(req: SummaryRequest, user: dict = Depends(get_current_user)):
    """Summarize using Gemini or fallback T5 with model info."""
    try:
        result = await summarize_text(req.text, req.mode)
        if isinstance(result, dict):
            summary = result.get("summary", "")
            model_used = result.get("model_used", "Unknown")
        else:
            summary = result
            model_used = "Gemini"
        return {"summary": summary, "model_used": model_used}
    except Exception as e:
        logger.error(f"Summarization failed: {e}")
        raise HTTPException(status_code=500, detail="Internal summarization error")

@app.post("/summarize-offline")
async def summarize_offline(request: dict, user: dict = Depends(get_current_user)):
    text = request.get("text", "")
    mode = request.get("mode", "short")
    summary = await summarize_with_t5(text, mode)
    return {"summary": summary, "model_used": "T5"}

@app.post("/quiz")
async def quiz(req: QuizRequest, user: dict = Depends(get_current_user)):
    quiz = await generate_quiz(req.text, req.difficulty)
    return {"quiz": quiz}

@app.post("/keywords")
async def keywords(req: TextRequest, user: dict = Depends(get_current_user)):
    keywords = await extract_keywords(req.text)
    return {"keywords": keywords}

# ------------------ File Summarization ------------------
async def summarize_file_text(text: str, mode: str):
    result = await summarize_text(text, mode)
    if isinstance(result, dict):
        return result.get("summary", ""), result.get("model_used", "Unknown")
    return result, "Gemini"

@app.post("/summarize-pdf")
async def summarize_pdf(file: UploadFile = File(...), mode: str = "short", user: dict = Depends(get_current_user)):
    try:
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        text = "".join([page.get_text("text") for page in doc])
        summary, model_used = await summarize_file_text(text, mode)
        return {"summary": summary, "model_used": model_used}
    except Exception as e:
        logger.error(f"PDF summarization failed: {e}")
        raise HTTPException(status_code=500, detail="Failed to summarize PDF")

@app.post("/summarize-docx")
async def summarize_docx(file: UploadFile = File(...), mode: str = "short", user: dict = Depends(get_current_user)):
    try:
        contents = await file.read()
        with open("temp.docx", "wb") as f:
            f.write(contents)
        doc = docx.Document("temp.docx")
        text = "\n".join([para.text for para in doc.paragraphs])
        summary, model_used = await summarize_file_text(text, mode)
        return {"summary": summary, "model_used": model_used}
    except Exception as e:
        logger.error(f"DOCX summarization failed: {e}")
        raise HTTPException(status_code=500, detail="Failed to summarize DOCX")

@app.post("/summarize-txt")
async def summarize_txt(file: UploadFile = File(...), mode: str = "short", user: dict = Depends(get_current_user)):
    try:
        contents = await file.read()
        text = contents.decode("utf-8")
        summary, model_used = await summarize_file_text(text, mode)
        return {"summary": summary, "model_used": model_used}
    except Exception as e:
        logger.error(f"TXT summarization failed: {e}")
        raise HTTPException(status_code=500, detail="Failed to summarize TXT")

@app.post("/summarize-pptx")
async def summarize_pptx(file: UploadFile = File(...), mode: str = "short", user: dict = Depends(get_current_user)):
    try:
        contents = await file.read()
        with open("temp.pptx", "wb") as f:
            f.write(contents)
        prs = Presentation("temp.pptx")
        text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        summary, model_used = await summarize_file_text(text, mode)
        return {"summary": summary, "model_used": model_used}
    except Exception as e:
        logger.error(f"PPTX summarization failed: {e}")
        raise HTTPException(status_code=500, detail="Failed to summarize PPTX")

# ------------------ Document Processing ------------------
@app.post("/process-pdf")
async def process_pdf(file: UploadFile = File(...), mode: str = "short", user: dict = Depends(get_current_user)):
    pdf_bytes = await file.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    text = "".join([page.get_text("text") for page in doc])
    summary, model_used = await summarize_file_text(text, mode)
    return {
        "summary": summary,
        "model_used": model_used,
        "quiz": await generate_quiz(text),
        "keywords": await extract_keywords(text),
        "flashcards": await generate_flashcards(text),
    }

@app.post("/process-docx")
async def process_docx(file: UploadFile = File(...), mode: str = "short", user: dict = Depends(get_current_user)):
    contents = await file.read()
    with open("temp.docx", "wb") as f:
        f.write(contents)
    doc = docx.Document("temp.docx")
    text = "\n".join([para.text for para in doc.paragraphs])
    summary, model_used = await summarize_file_text(text, mode)
    return {
        "summary": summary,
        "model_used": model_used,
        "quiz": await generate_quiz(text),
        "keywords": await extract_keywords(text),
        "flashcards": await generate_flashcards(text),
    }

@app.post("/process-txt")
async def process_txt(file: UploadFile = File(...), mode: str = "short", user: dict = Depends(get_current_user)):
    contents = await file.read()
    text = contents.decode("utf-8")
    summary, model_used = await summarize_file_text(text, mode)
    return {
        "summary": summary,
        "model_used": model_used,
        "quiz": await generate_quiz(text),
        "keywords": await extract_keywords(text),
        "flashcards": await generate_flashcards(text),
    }

@app.post("/process-pptx")
async def process_pptx(file: UploadFile = File(...), mode: str = "short", user: dict = Depends(get_current_user)):
    contents = await file.read()
    with open("temp.pptx", "wb") as f:
        f.write(contents)
    prs = Presentation("temp.pptx")
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    summary, model_used = await summarize_file_text(text, mode)
    return {
        "summary": summary,
        "model_used": model_used,
        "quiz": await generate_quiz(text),
        "keywords": await extract_keywords(text),
        "flashcards": await generate_flashcards(text),
    }

# ------------------ CORS ------------------
from fastapi.middleware.cors import CORSMiddleware

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

